#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""套模版风格构建 TRB 0701 组会汇报 PPT。中文文字·无破折号·无双引号·英文图表。"""
import glob, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

os.chdir('/Users/weiyutang/Desktop/TRB/Paper/0701组会汇报')
TMPL = glob.glob('*.pptx')[0]
RES = '/Users/weiyutang/Desktop/TRB/结果/0628-追越场景审核'
prs = Presentation(TMPL)
EMU_W, EMU_H = prs.slide_width, prs.slide_height
SW, SH = EMU_W/914400, EMU_H/914400   # inch

# 模版配色
ACCENT = RGBColor(0x15, 0x60, 0x82)   # 156082
NAVY   = RGBColor(0x0E, 0x28, 0x41)
DARK   = RGBColor(0x22, 0x22, 0x22)
GRAY   = RGBColor(0x55, 0x55, 0x55)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)   # 浅冰蓝, 深底副标题用
BLACK  = RGBColor(0x00, 0x00, 0x00)   # 一级标题(除封面)用黑色
CN = '微软雅黑'

# 版式索引 by name
LAY = {l.name: l for l in prs.slide_layouts}
BLANK = LAY.get('页面')  # 用页面版式(带标题占位)
SECTION = LAY.get('章节标题')

# 清理封面版式里烤死的模版残留文字与胶囊(只清幻灯片占位符清不掉这些 layout 级形状),
# 否则新标题会与 — 2025-2026学年船舶海洋与建筑工程学院 — 等旧文字重叠。封面与致谢共用此版式。
_LEFTOVER = {'文本框 41', '文本框 44', '直接连接符 11', '文本占位符 8'}
for _sh in list(LAY['封面'].shapes):
    if _sh.name in _LEFTOVER:
        _sh._element.getparent().remove(_sh._element)

def _set_cn(run, font=CN):
    run.font.name = font
    from pptx.oxml.ns import qn
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {}); rpr.append(ea)
    ea.set('typeface', font)

def add_text(slide, l, t, w, h, text, size=16, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=CN, sp_after=6):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        # 支持行首 * = 加粗小标题
        b = bold
        if ln.startswith('* '):
            ln = ln[2:]; b = True
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = color; _set_cn(r, font)
    return tb

def add_title(slide, text):
    return add_text(slide, 0.55, 0.28, SW-1.1, 0.9, text, size=26, color=BLACK, bold=True)

def add_img(slide, path, l, t, maxw, maxh, center_in=True):
    iw, ih = Image.open(path).size
    ar = iw/ih
    w = maxw; h = w/ar
    if h > maxh:
        h = maxh; w = h*ar
    lx = l + (maxw-w)/2 if center_in else l
    ty = t + (maxh-h)/2 if center_in else t
    slide.shapes.add_picture(path, Inches(lx), Inches(ty), Inches(w), Inches(h))

def new_slide(title=None, section=False):
    s = prs.slides.add_slide(SECTION if section else BLANK)
    # 清空版式占位里的示例文字
    for ph in list(s.placeholders):
        try: ph.text_frame.clear()
        except Exception: pass
    if title and not section:
        add_title(s, title)
    return s

# ================= 删除模版全部页(drop_rel 防孤儿 slide 重名) =================
from pptx.oxml.ns import qn
sldIdLst = prs.slides._sldIdLst
for sid in list(sldIdLst):
    rId = sid.get(qn('r:id'))
    try: prs.part.drop_rel(rId)
    except Exception: pass
    sldIdLst.remove(sid)

# ================= 封面(从封面版式重建) =================
cover = prs.slides.add_slide(LAY['封面'])
for ph in list(cover.placeholders):
    try: ph.text_frame.clear()
    except Exception: pass
# 标题放在下半浅色线稿区, 用深色(NAVY/ACCENT/DARK)保证在浅底上对比, 避开上方校名书法
add_text(cover, 1.0, 3.9, SW-2.0, 1.4, '连续动作空间下可证明合规的\n无人船智能避碰',
         size=32, color=NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(cover, 1.0, 5.33, SW-2.0, 0.5, '投影式安全盾 与 强化学习', size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(cover, 1.0, 6.55, SW-2.0, 0.5, '汇报人  唐伟昱      指导教师  薛杰      2026.07.01 组会汇报',
         size=14, color=DARK, align=PP_ALIGN.CENTER)

# ================= 内容页 =================
# S2 研究背景与动机
s = new_slide('一、研究背景与动机')
add_text(s, 0.55, 1.3, SW-1.1, 5.4,
 '* 场景\n无人船在开阔水域航行, 必须遵守国际海上避碰规则 COLREGs, 例如对遇时双方向右转让、交叉相遇时给路船避让。\n\n'
 '* 需求\n用强化学习训练避碰策略时, 单靠奖励只能鼓励合规, 不能保证合规。安全攸关的海事场景需要的是可以证明的规则遵守, 而不是大概率遵守。\n\n'
 '* 意义\n把可以证明合规的安全机制引入连续控制的无人船避碰, 既保住合规硬保证, 又保留连续控制的平滑与灵活。',
 size=16)

# S3 问题与挑战
s = new_slide('二、问题与挑战')
add_text(s, 0.55, 1.3, 5.6, 5.2,
 '现有工作存在一个两难:\n\n'
 '* 路线一: 可以证明合规\n把动作离散成有限个, 用动作屏蔽保证每步都合规。代价是放弃连续控制, 动作生硬。\n\n'
 '* 路线二: 连续控制\n直接在连续动作上训练, 保留平滑灵活。代价是合规只能靠软奖励, 没有硬保证。\n\n'
 '* 我们要打通的空白\n在连续动作空间上同时拿到可以证明的合规和连续控制。', size=16)
# 右侧对比示意(用两栏色块文字)
add_text(s, 6.5, 1.6, 6.2, 1.2, '可证明合规', size=18, color=ACCENT, bold=True)
add_text(s, 6.5, 2.05, 6.2, 1.0, '只能离散动作 (Krasowski 2024)', size=15, color=GRAY)
add_text(s, 6.5, 3.4, 6.2, 1.2, '连续控制', size=18, color=ORANGE, bold=True)
add_text(s, 6.5, 3.85, 6.2, 1.0, '只能软奖励 · 无硬保证', size=15, color=GRAY)
add_text(s, 6.5, 5.2, 6.2, 1.3, '本文: 连续动作 + 可证明合规', size=18, color=NAVY, bold=True)

# S4 数据集与实验设置
s = new_slide('三、数据集与实验设置')
add_text(s, 0.55, 1.3, SW-1.1, 5.2,
 '* 仿真环境\nCommonOcean 开源海事场景, 船舶为偏航受限动力学模型, 连续控制输入为加速度与转艏率。\n\n'
 '* 会遇类型\n本阶段聚焦对遇和交叉两类标准会遇。追越类型的场景已经构建完成, 留待后续训练。\n\n'
 '* 数据划分 (冻结不变)\n训练 94 个场景, 测试 40 个从不参与训练的场景。\n\n'
 '* 评估口径\n多个随机种子, 报告到达率的中位数与四分位均值, 以及失败率, 全部指标同口径对比。', size=16)

# S4b COLREGs 会遇分类与合规角色
s = new_slide('三、COLREGs 会遇分类与合规角色')
add_img(s, 'fig_encounter.png', 0.4, 1.35, SW-0.8, 3.3)
add_text(s, 0.55, 4.9, SW-1.1, 2.1,
 '* 四类标准会遇\n状态机把两船几何判为对遇, 交叉, 追越, 直航之一, 再加无冲突与紧急两种, 共六种态势。\n'
 '* 让路船与直航船\n每类会遇里本船被指派角色, 让路船须做出明显避让, 对遇与交叉一律向右转, 直航船须保向保速。\n'
 '* 条款号纪律\n汇报一律用国际规则真实条款号, 对遇 Rule 14, 交叉 Rule 15, 追越 Rule 13, 直航 Rule 17。', size=13, sp_after=4)

# S4c 研究对象与船舶参数
s = new_slide('三、研究对象与船舶参数')
add_img(s, 'fig_ship.png', 0.5, 1.5, SW-1.0, 2.6)
add_text(s, 0.55, 4.35, SW-1.1, 2.4,
 '* 船型与模型\n本船为集装箱船 SR108, 采用 Krasowski 偏航受限点质量模型, 忠实沿用不改。\n'
 '* 连续控制是卖点锚点\n智能体直接输出加速度与转艏率两个连续量, 这是与离散动作方法的根本区别。\n'
 '* 碰撞判定\n用 175 乘 25.4 米的旋转矩形占据判交叠, 不是简化的圆形半径。', size=13, sp_after=4)

# S4d 观测空间
s = new_slide('三、观测空间: 策略看见什么')
add_img(s, 'fig_obs_space.png', 0.4, 1.45, SW-0.8, 3.4)
add_text(s, 0.55, 5.25, SW-1.1, 1.6,
 '策略输入是 27 维实数向量, 逐字复现 Krasowski, 分为本船自身, 目标, 四个方向扇区的他船, 以及五个终止标志四组。这是理解策略在学什么的前提, 观测口径与离散基线完全一致以保证对比干净。', size=13, color=DARK)

# S4e 到达判定与终止条件
s = new_slide('三、到达判定与终止条件')
add_text(s, 0.55, 1.2, SW-1.1, 0.5, '到达不是简单的到点, 采用 CommonOcean 官方口径, 三条同时成立才算到达', size=15, color=GRAY, bold=True)
add_text(s, 0.95, 1.95, SW-1.9, 1.5,
 '* 位置进入目标区\n* 艏向落在目标朝向区间\n* 时间步落在目标时间区间', size=15, color=NAVY, sp_after=3)
add_text(s, 0.55, 3.85, SW-1.1, 3.0,
 '* 目标区\n约 400 乘 60 米的矩形, 距起点约 4500 米。\n'
 '* 五个终止条件, 任一成立即结束一局\n到达, 碰撞(船体多边形相交), 超时(最多 170 步, 每步 10 秒), 出界, 停船。\n'
 '* 为什么强调\n到达要求朝向达标且按时进入, 比单纯到点更严, 这直接影响到达率数字的解读。', size=13, sp_after=4)

# S4f 训练设置与评估指标定义
s = new_slide('三、训练设置与评估指标定义')
add_text(s, 0.55, 1.25, 6.15, 5.4,
 '* 训练算法\n连续主线与离散基线都用 PPO, 区别只在于连续臂用投影盾, 离散臂用动作屏蔽, 网络均为两层各 64。\n'
 '* 数据与种子 (与结果页同口径)\n基准为双船相遇场景, 主对比在冻结的对遇与交叉子集上, 94 个训练, 40 个测试, 5 个随机种子。\n'
 '* 训练协议\n折扣 0.99, 网络与训练步数忠实沿用 Krasowski, 投影盾在环境内执行。', size=13, sp_after=4)
add_text(s, 6.95, 1.25, 5.95, 5.4,
 '* 评估指标定义\n到达率, 成功进入目标区的局数比例。\n碰撞率, 发生船体碰撞的局数比例, 带盾方法全程为零。\n违规率, 直航态势按步计, 让路态势按整段相遇计, 与 Krasowski 同口径。\n四分位均值 IQM, 去掉最高最低各四分之一种子后的均值, 抗离群更稳。\n控制 jerk, 相邻两步控制量的变化幅度, 越低越平滑。', size=13, sp_after=4)

# S5 相关工作与对比基线
s = new_slide('四、相关工作与对比基线')
add_text(s, 0.55, 1.3, SW-1.1, 4.0,
 '* Krasowski 与 Althoff 2024 (对比基线)\n把 COLREGs 形式化为状态机, 用离散动作屏蔽实现可证明合规。我们在自己的代码里独立复现这一基线用于对比。\n\n'
 '* Markgraf 等 2026 (方法学基础)\n投影式安全盾的理论与开源实现, 发表于机器学习期刊 TMLR。我们借用其投影盾思路。\n\n'
 '* 我们的创新点\n首次把投影式安全盾落到海事 COLREGs 的连续动作空间, 正面解决离散方法自身承认的局限。', size=16)

# S5b 安全盾放在策略里还是环境里 (SE-RL vs SP-RL)
s = new_slide('四、安全盾放在策略里还是环境里')
add_img(s, 'fig_serl_sprl.png', 0.5, 1.4, SW-1.0, 3.9)
add_text(s, 0.55, 5.5, SW-1.1, 1.5,
 '* 两种接线方式\n安全盾可以放进策略里参与梯度回传, 也可以放进环境里作为独立一层。\n'
 '* 本文选择盾即环境\n盾对策略透明, 训练与部署用同一套投影, 实现简单且与离散基线可比。安全强化学习常见三条路线是动作替换, 投影, 动作屏蔽, 本文走连续投影这条。', size=13, sp_after=4)

# S6 方法框架图
s = new_slide('五、方法总览: 投影式安全盾框架')
add_img(s, 'fig_framework.png', 0.5, 1.3, SW-1.0, 4.6)
add_text(s, 0.55, 6.05, SW-1.1, 1.0,
 '策略输出期望动作后, 状态机判定当前会遇态势与合规方向, 构造出同时满足合规方向与无碰撞的动作约束集, 再用二次规划把期望动作投影进该集合, 得到安全动作执行。', size=13.5, color=GRAY)

# S7 方法·第一类固定基础
s = new_slide('六、方法 (一): 固定不动的基础')
add_text(s, 0.55, 1.15, SW-1.1, 0.5, '这些是我们不改动的平台, 保证对比干净可比', size=15, color=GRAY, bold=True)
add_text(s, 0.55, 1.75, SW-1.1, 5.0,
 '* 仿真环境与船舶动力学\n沿用 CommonOcean 与偏航受限模型, 不改。\n\n'
 '* COLREGs 形式化状态机\n复现 Krasowski 的相遇态势分类与合规方向判定, 不改。\n\n'
 '* 投影式安全盾\n核心创新架构, 把合规与无碰撞刻画为约束集并投影, 固定不动。\n\n'
 '* 基础奖励函数\n忠实复现 Krasowski 原始系数, 不改基础系数。\n\n'
 '* 训练算法与离散对比基线\n连续 PPO 固定; 离散基线忠实复现, 不加论文里没有的东西。', size=15.5)

# S8 方法·基础奖励数学
s = new_slide('六、方法 (一): 基础奖励函数')
add_img(s, 'fig_math_reward.png', 0.5, 1.4, SW-1.0, 4.0)
add_text(s, 0.55, 5.7, SW-1.1, 1.2,
 '总奖励由五部分组成, 系数全部来自 Krasowski 原文。其中接近奖励是随距离缩短给出的稠密引导, 进入目标只在到达那一步给一次奖励。', size=13.5, color=GRAY)

# S8b 方法·投影盾数学机制(第一类核心·公式照 方法数学公式.md §2)
s = new_slide('六、方法 (一): 投影式安全盾的数学机制')
add_img(s, 'fig_math_shield.png', 0.5, 1.12, SW-1.0, 2.15)
add_text(s, 0.55, 3.4, SW-1.1, 3.0,
 '* 状态机判合规方向\n给路船(对遇, 右舷交叉)被约束为向右转; 直航船被约束为保向保速, 转艏率与加速度都压在很小的窄带内; 无冲突态势不施加方向约束。\n\n'
 '* 无碰撞约束\n要求本船下一步的占据区域与每艘他船的预测占据区域不相交, 并留出安全裕度, 线性化后得到一组分离超平面约束。\n\n'
 '* 二次规划投影与兜底\n用二次规划把策略输出的期望控制以最小改动投影进既合规又无碰撞的可行集; 可行集为空或出现紧急态势时, 放松合规以保证无碰撞, 或切换到紧急控制器。',
 size=13, sp_after=4)
add_text(s, 0.55, 6.5, SW-1.1, 0.6,
 '诚实定位: 方向合规由投影约束构造性保证(命题二, 对所有种子成立); 零船体碰撞目前为经验结果(档位A), 完整证明是后续工作。', size=12, color=GRAY)

# S9 方法·第二类我们调的(文上·两类塑形数学图下·well_B 已保留 + c_step 主攻)
s = new_slide('七、方法 (二): 我们在调整的方法')
add_text(s, 0.55, 1.15, SW-1.1, 2.5,
 '目标是在第一类固定基础之上叠加塑形, 让连续策略学得更好, 环境, 状态机与安全盾都不改动。\n\n'
 '* 进门势塑形 well_B (已保留)\n给近目标区一个稠密梯度帮助学习进入, 把种子间方差砍半, 属于策略不变的势函数塑形。\n'
 '* 每步生存成本 c_step (当前主攻, 标定中)\n每步给一点微小成本, 让在目标附近往复而不进入变得不划算, 从根本上改变奖励地形。\n'
 '* 动作平滑 (计划)\n惩罚相邻两步控制的变化, 降低抖动。', size=13.5, sp_after=4)
add_img(s, 'fig_math_shaping.png', 0.5, 3.95, SW-1.0, 2.9)

# S9b 可证明的部分与诚实边界 (命题1/命题2/档位A-B, 借鉴 0626 论文)
s = new_slide('七、可证明的部分与诚实边界')
add_text(s, 0.55, 1.3, SW-1.1, 5.4,
 '* 命题一, 塑形不改最优\n进门势 well_B 是纯状态的势函数塑形, 按 Ng 1999 的策略不变性定理不改变最优策略, 因此不破坏安全盾的任何约束与合规保证。这正是它区别于一般软奖励的关键。\n\n'
 '* 命题二, 合规方向的构造性保证\n只要安全集非空且投影有解, 安全动作就一定落在合规方向集里, 让路时向右转, 直航时保向保速。这个保证是构造出来的, 对所有随机种子都成立, 与策略学得好不好无关。例外是紧急兜底那一步, 方向约束被放松以优先避碰, 这一步的比例我们如实报告。\n\n'
 '* 诚实边界, 当前是档位A\n零船体碰撞目前是经验结果, 来自单前瞻的一阶近似加保守裕度, 我们不声称可证明零碰撞。把递归可行性做进硬约束以拿到可证明零碰撞是档位B, 属于后续工作, 它决定英文标题能不能用可证明无碰撞这个说法。', size=13.5, sp_after=4)

# S10 实验结果·主对比
s = new_slide('八、实验结果: 主对比 (同一冻结数据集)')
add_img(s, 'table_main_comparison.png', 0.5, 1.25, SW-1.0, 2.2)
add_img(s, 'fig_training_curves.png', 0.5, 3.6, SW-1.0, 2.7)
add_text(s, 0.55, 6.45, SW-1.1, 0.7,
 '全程零碰撞是安全盾的核心贡献且稳定。当前离散基线在到达率与控制平滑上仍占优, 连续方法的平滑卖点待后续工作。', size=12.5, color=GRAY)

# S11 训练稳定性
s = new_slide('八、实验结果: 训练稳定性与每步生存成本')
add_img(s, 'fig_seed_recovery.png', 0.5, 1.3, SW-1.0, 3.5)
add_img(s, 'table_cstep.png', 3.4, 5.0, 6.4, 1.9)
add_text(s, 0.55, 1.05, SW-1.1, 0.35, '部分种子到达率为零的原因是进入目标的奖励太弱; 每步生存成本随取值增大单调恢复失败种子', size=12.5, color=GRAY)

# S12 深入可视化
s = new_slide('九、深入可视化: 轨迹与控制平滑度')
add_img(s, 'fig_trajectory_control.png', 0.5, 1.3, SW-1.0, 3.0)
add_img(s, 'fig_dither_preview.png', 0.5, 4.5, 6.0, 2.4, center_in=False)
add_text(s, 6.7, 4.7, 6.1, 2.2,
 '诚实定位\n\n'
 '当前无平滑处理的连续控制在方向盘轴上比离散更抖。\n\n'
 '动作平滑处理已在历史数据上验证, 可把抖动降约三分之一且到达率不掉, 是下一步要施加的改进。', size=14, color=DARK)

# S13 已完成·追越场景
s = new_slide('九、已完成工作: 追越场景构建')
add_text(s, 0.55, 1.1, SW-1.1, 0.5, '追越会遇的场景已经构建并做质量审核, 数据侧已铺好路, 追越是待训练而非待构建', size=15, color=GRAY, bold=True)
imgs = [f'{RES}/1_追越_全覆盖演示_9类.png', f'{RES}/4_最经典超车_同一条线.png', f'{RES}/2_追越_近景放大.png']
xs = [0.5, 4.9, 9.0]
for p, x in zip(imgs, xs):
    if os.path.exists(p): add_img(s, p, x, 1.9, 3.9, 4.4, center_in=False)

# S14 失败方案时间线
s = new_slide('十、失败方案时间线')
add_text(s, 0.55, 1.25, SW-1.1, 5.3,
 '* 换用 SAC 作主线算法\n与投影盾信用分配冲突, 梯度发散, 转为 PPO。\n'
 '* 增加训练数据\n崩塌率不变, 与数据量无关。\n'
 '* 调节探索熵\n调低导致崩塌, 调高导致到达率下降, 两个方向都伤。\n'
 '* 动作混叠惩罚\n过强, 多数种子崩塌。\n'
 '* 横向进带势\n引发灾难性崩塌, 净收益为负。\n'
 '* 终端保速势\n只是重新洗牌, 没有净收益。\n'
 '* 距离课程学习\n场景距离基本恒定, 课程没有梯度可用。\n'
 '* 热启动与路径惩罚\n前者是复用已有模型并非真鲁棒, 后者打错目标且产生退化解。', size=14.5, sp_after=3)

# S15 待解决与下一步
s = new_slide('十一、待解决的问题与下一步')
add_text(s, 0.55, 1.35, SW-1.1, 5.0,
 '* 标定每步生存成本的取值\n找到既能恢复失败种子, 又不影响健康种子的取值。\n\n'
 '* 平滑度卖点\n施加动作平滑, 并用油门轴和物理空间指标证明连续控制更平滑且可视化。\n\n'
 '* 四方对比大图\n补齐无合规和软奖励两个基线, 完成完整对比。\n\n'
 '* 攻克追越会遇\n在已经构建的追越场景上训练连续策略。', size=16)

# ================= 致谢页(封面版式) =================
th = prs.slides.add_slide(LAY['封面'])
for ph in list(th.placeholders):
    try: ph.text_frame.clear()
    except Exception: pass
add_text(th, 1.0, 4.05, SW-2.0, 1.2, '感谢聆听  请批评指正', size=32, color=NAVY, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(th, 1.0, 5.5, SW-2.0, 0.5, '唐伟昱      2026.07.01', size=15, color=DARK, align=PP_ALIGN.CENTER)

# ================= 保存 =================
out = '/Users/weiyutang/Desktop/TRB/Paper/0701组会汇报/TRB_0701组会汇报.pptx'
prs.save(out)
print('✅ 保存', out, '| 总页数', len(prs.slides))
